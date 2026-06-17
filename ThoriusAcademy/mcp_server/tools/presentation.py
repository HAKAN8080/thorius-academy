"""
PowerPoint builder for Thorius Academy.

Slides are described declaratively in JSON. Supported layouts:
  - cover                 : full-bleed navy title slide w/ logo
  - section_divider       : full-navy section break with copper accent
  - content_typography    : numbered-list editorial layout (e.g. learning outcomes)
  - three_columns         : 3 vertical persona/role/comparison columns
  - timeline              : horizontal zigzag timeline
  - numbered_grid         : grid of bordered cards with kicker labels
  - kpi_tiles             : 2x2 or 4-up big-number tiles
  - formula_with_example  : highlighted formula + worked example panel
  - table                 : branded data table with intro + footer note
  - bar_chart, line_chart : data viz with side notes
  - bullets               : prose bullet summary (for section summaries)
  - closing               : full-navy thank-you slide
"""

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from helpers import ensure_section, get_section_dir, mark_step

# ============================================================
# BRAND CONSTANTS — Thorius Academy
# ============================================================
NAVY      = RGBColor(0x1A, 0x2B, 0x4A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2A, 0x2A, 0x2A)
GREY_MID  = RGBColor(0x6B, 0x6B, 0x6B)
GREY_LITE = RGBColor(0xE5, 0xE5, 0xE5)
COPPER    = RGBColor(0xB8, 0x96, 0x5A)

BRAND_DIR = Path(__file__).parent.parent / "brand" / "assets"
LOGO_PATH = BRAND_DIR / "logo_transparent.png"
BG_PATH = BRAND_DIR / "bg_pattern_light.png"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
HEADER_H = Inches(0.85)
FOOTER_H = Inches(0.4)
FOOTER_TOP = Inches(7.0)


# ============================================================
# LOW-LEVEL HELPERS
# ============================================================
def add_rect(slide, x, y, w, h, fill, line=None, line_width=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_width:
            s.line.width = line_width
    return s


def add_text(slide, x, y, w, h, text, size=18, bold=False, italic=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
    return tb


def add_hline(slide, x1, y, x2, color=GREY_LITE, width_pt=0.5):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y, x2, y)
    ln.line.color.rgb = color
    ln.line.width = Pt(width_pt)
    return ln


def add_vline(slide, x, y1, y2, color=GREY_LITE, width_pt=0.5):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y1, x, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(width_pt)
    return ln


def add_body_bg(slide):
    if not BG_PATH.exists():
        return
    try:
        bg = slide.shapes.add_picture(
            str(BG_PATH), 0, HEADER_H,
            SLIDE_W, FOOTER_TOP - HEADER_H
        )
        spTree = bg._element.getparent()
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)
    except Exception:
        pass


def add_header(slide, title, kicker=None):
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, HEADER_H)
    h.fill.solid()
    h.fill.fore_color.rgb = NAVY
    h.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), 0, Inches(9.5), HEADER_H)
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    for r in p.runs:
        r.font.name = "Calibri"
        r.font.size = Pt(26)
        r.font.bold = True
        r.font.color.rgb = WHITE
    if kicker:
        tb2 = slide.shapes.add_textbox(SLIDE_W - Inches(3.5), 0,
                                        Inches(3.0), HEADER_H)
        tf2 = tb2.text_frame
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf2.paragraphs[0]
        p.text = kicker
        p.alignment = PP_ALIGN.RIGHT
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.size = Pt(11)
            r.font.bold = True
            r.font.color.rgb = COPPER


def add_footer(slide, page_no, total):
    y = SLIDE_H - FOOTER_H - Inches(0.05)
    add_hline(slide, Inches(0.4), y, SLIDE_W - Inches(0.4),
              color=GREY_LITE, width_pt=0.5)
    tb_l = slide.shapes.add_textbox(Inches(0.4), SLIDE_H - FOOTER_H,
                                     Inches(6), FOOTER_H)
    p = tb_l.text_frame.paragraphs[0]
    p.text = "Thorius Academy"
    for r in p.runs:
        r.font.name = "Calibri"
        r.font.size = Pt(10)
        r.font.color.rgb = GREY_MID
    tb_r = slide.shapes.add_textbox(SLIDE_W - Inches(2.4),
                                     SLIDE_H - FOOTER_H,
                                     Inches(2), FOOTER_H)
    p = tb_r.text_frame.paragraphs[0]
    p.text = f"Page {page_no} / {total}"
    p.alignment = PP_ALIGN.RIGHT
    for r in p.runs:
        r.font.name = "Calibri"
        r.font.size = Pt(10)
        r.font.color.rgb = GREY_MID


# ============================================================
# LAYOUT BUILDERS
# ============================================================

def layout_cover(slide, spec):
    """
    Full navy cover with logo, title, subtitle, brand tag.
    spec keys:
      title (str, required)
      subtitle (str)
      kicker (str)    — small copper label above title (default: brand)
      tagline (str)   — small italic line below brand
      meta (str)      — bottom-right meta line
    """
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    # Left vertical copper bar
    add_rect(slide, Inches(0.8), Inches(1.6), Inches(0.08), Inches(4.3), COPPER)

    kicker = spec.get("kicker", "THORIUS ACADEMY  •  COURSE")
    add_text(slide, Inches(1.1), Inches(1.55), Inches(11), Inches(0.4),
             kicker, size=12, bold=True, color=COPPER)

    title = spec.get("title", "Untitled")
    add_text(slide, Inches(1.1), Inches(2.0), Inches(11.0), Inches(1.6),
             title, size=54, bold=True, color=WHITE)

    subtitle = spec.get("subtitle", "")
    if subtitle:
        add_text(slide, Inches(1.1), Inches(4.4), Inches(11), Inches(0.6),
                 subtitle, size=22, italic=True, color=COPPER)

    add_hline(slide, Inches(1.1), Inches(5.3), Inches(7.5),
              color=WHITE, width_pt=1.0)

    add_text(slide, Inches(1.1), Inches(5.85), Inches(7), Inches(0.4),
             "A THORIUS ACADEMY COURSE",
             size=11, bold=True, color=COPPER)
    tagline = spec.get("tagline",
                        "Practical frameworks for retail teams")
    add_text(slide, Inches(1.1), Inches(6.25), Inches(9), Inches(0.4),
             tagline, size=13, italic=True, color=GREY_LITE)

    if LOGO_PATH.exists():
        try:
            logo_w = Inches(2.6)
            logo_h = Inches(1.45)
            slide.shapes.add_picture(
                str(LOGO_PATH),
                SLIDE_W - logo_w - Inches(1.0), Inches(0.9),
                logo_w, logo_h
            )
        except Exception:
            pass

    meta = spec.get("meta", "")
    if meta:
        add_text(slide, SLIDE_W - Inches(4.5), SLIDE_H - Inches(0.75),
                 Inches(4.0), Inches(0.35),
                 meta, size=11, color=GREY_LITE, align=PP_ALIGN.RIGHT)


def layout_section_divider(slide, spec):
    """Full navy. spec: section_no, title."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    section_no = spec.get("section_no", "")
    add_text(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(0.6),
             f"SECTION {section_no}", size=18, bold=True,
             color=COPPER, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(1.2),
             spec.get("title", ""), size=36, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    ln = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        (SLIDE_W - Inches(2)) / 2, Inches(4.5),
        Inches(2), Emu(20000)
    )
    ln.fill.solid()
    ln.fill.fore_color.rgb = COPPER
    ln.line.fill.background()


def layout_content_typography(slide, spec):
    """
    Numbered editorial layout (used for learning outcomes etc).
    Two columns × N rows. Each item: big copper number | thin rule | title + body.
    spec:
      intro (str, italic)
      items (list of {number, title, body})
    """
    intro = spec.get("intro", "")
    if intro:
        add_text(slide, Inches(0.6), Inches(1.15), Inches(12.2), Inches(0.5),
                 intro, size=17, italic=True, color=NAVY)

    items = spec.get("items", [])
    col_w = Inches(6.05)
    row_h = Inches(1.45)
    start_x = Inches(0.6)
    start_y = Inches(2.05)
    gap_x = Inches(0.15)
    gap_y = Inches(0.15)
    for idx, item in enumerate(items):
        col = idx % 2
        row = idx // 2
        x = start_x + (col_w + gap_x) * col
        y = start_y + (row_h + gap_y) * row
        add_text(slide, x, y - Inches(0.05), Inches(1.2), row_h + Inches(0.1),
                 item.get("number", f"{idx+1:02d}"),
                 size=48, bold=True, color=COPPER)
        add_vline(slide, x + Inches(1.3), y + Inches(0.15),
                  y + row_h - Inches(0.15),
                  color=GREY_LITE, width_pt=0.75)
        add_text(slide, x + Inches(1.45), y + Inches(0.05),
                 col_w - Inches(1.5), Inches(0.45),
                 item.get("title", ""), size=16, bold=True, color=NAVY)
        add_text(slide, x + Inches(1.45), y + Inches(0.55),
                 col_w - Inches(1.5), row_h - Inches(0.55),
                 item.get("body", ""), size=12, color=TEXT_DARK)


def layout_three_columns(slide, spec):
    """
    Three vertical persona/role columns separated by thin rules.
    spec:
      intro (str)
      columns (list of {label, title, lead, items[]})
      footer_label (str)
      footer_text (str)
    """
    intro = spec.get("intro", "")
    if intro:
        add_text(slide, Inches(0.6), Inches(1.15), Inches(12.2), Inches(0.5),
                 intro, size=17, italic=True, color=NAVY)

    columns = spec.get("columns", [])[:3]
    col_w = Inches(4.0)
    start_x = Inches(0.6)
    y_top = Inches(2.1)
    gap = Inches(0.13)
    for i, p in enumerate(columns):
        x = start_x + (col_w + gap) * i
        if i > 0:
            add_vline(slide, x - gap / 2, y_top, Inches(6.2),
                      color=GREY_LITE, width_pt=0.75)
        add_text(slide, x, y_top, col_w, Inches(0.3),
                 p.get("label", ""), size=11, bold=True, color=COPPER)
        add_text(slide, x, y_top + Inches(0.3), col_w, Inches(0.6),
                 p.get("title", ""), size=28, bold=True, color=NAVY)
        add_rect(slide, x, y_top + Inches(0.95),
                 Inches(0.5), Inches(0.04), COPPER)
        add_text(slide, x, y_top + Inches(1.15), col_w, Inches(0.6),
                 p.get("lead", ""), size=14, italic=True, color=GREY_MID)
        for j, item in enumerate(p.get("items", [])):
            ty = y_top + Inches(1.85) + Inches(0.45) * j
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                          x + Inches(0.02),
                                          ty + Inches(0.13),
                                          Inches(0.12), Inches(0.12))
            dot.fill.solid()
            dot.fill.fore_color.rgb = NAVY
            dot.line.fill.background()
            add_text(slide, x + Inches(0.25), ty,
                     col_w - Inches(0.25), Inches(0.4),
                     item, size=13, color=TEXT_DARK)

    footer_label = spec.get("footer_label", "")
    footer_text = spec.get("footer_text", "")
    if footer_label or footer_text:
        band_y = Inches(6.3)
        add_hline(slide, Inches(0.6), band_y, Inches(12.7),
                  color=COPPER, width_pt=1.0)
        if footer_label:
            add_text(slide, Inches(0.6), band_y + Inches(0.1),
                     Inches(12.1), Inches(0.4),
                     footer_label, size=10, bold=True, color=COPPER)
        if footer_text:
            add_text(slide, Inches(0.6), band_y + Inches(0.4),
                     Inches(12.1), Inches(0.4),
                     footer_text, size=13, color=TEXT_DARK)


def layout_timeline(slide, spec):
    """
    Horizontal timeline with alternating top/bottom captions.
    spec:
      intro (str)
      items (list of {num, title, sub})
      footer (str)
    """
    intro = spec.get("intro", "")
    if intro:
        add_text(slide, Inches(0.6), Inches(1.15), Inches(12.2), Inches(0.5),
                 intro, size=17, italic=True, color=NAVY)

    items = spec.get("items", [])
    if not items:
        return

    rail_y = Inches(4.1)
    add_hline(slide, Inches(0.8), rail_y, SLIDE_W - Inches(0.8),
              color=NAVY, width_pt=2.5)

    n = len(items)
    span_left = Inches(0.8)
    span_right = SLIDE_W - Inches(0.8)
    step = (span_right - span_left) / (n - 1) if n > 1 else 0

    for i, item in enumerate(items):
        cx = span_left + step * i
        is_top = (i % 2 == 0)
        dot_d = Inches(0.32)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      cx - dot_d / 2,
                                      rail_y - dot_d / 2,
                                      dot_d, dot_d)
        dot.fill.solid()
        dot.fill.fore_color.rgb = COPPER
        dot.line.color.rgb = WHITE
        dot.line.width = Pt(2)

        cap_w = Inches(1.55)
        cap_h = Inches(1.45)
        if is_top:
            cap_y = rail_y - Inches(1.7)
            add_vline(slide, cx, rail_y - Inches(0.15), cap_y + cap_h,
                      color=GREY_LITE, width_pt=0.75)
        else:
            cap_y = rail_y + Inches(0.25)
            add_vline(slide, cx, rail_y + Inches(0.15), cap_y,
                      color=GREY_LITE, width_pt=0.75)

        add_text(slide, cx - cap_w / 2, cap_y, cap_w, Inches(0.5),
                 item.get("num", str(i + 1)),
                 size=22, bold=True, color=COPPER, align=PP_ALIGN.CENTER)
        add_text(slide, cx - cap_w / 2, cap_y + Inches(0.5),
                 cap_w, Inches(0.4),
                 item.get("title", ""),
                 size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, cx - cap_w / 2, cap_y + Inches(0.9),
                 cap_w, Inches(0.5),
                 item.get("sub", ""),
                 size=10, color=GREY_MID, align=PP_ALIGN.CENTER)

    footer = spec.get("footer", "")
    if footer:
        add_text(slide, Inches(0.6), Inches(6.4), Inches(12.2), Inches(0.4),
                 footer, size=13, bold=True, color=COPPER,
                 align=PP_ALIGN.CENTER, italic=True)


def layout_bullets(slide, spec):
    """Simple bulleted summary list.
    spec: items (list of str), bullet_size (int, default 17)
    """
    items = spec.get("items", [])
    size = spec.get("bullet_size", 17)
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.5),
                                   Inches(11.9), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.space_after = Pt(10)
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.size = Pt(size)
            r.font.color.rgb = TEXT_DARK


def layout_closing(slide, spec):
    """Full navy thank-you slide."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    if LOGO_PATH.exists():
        try:
            logo_w = Inches(2.8)
            logo_h = Inches(1.55)
            slide.shapes.add_picture(
                str(LOGO_PATH),
                (SLIDE_W - logo_w) / 2, Inches(1.3),
                logo_w, logo_h
            )
        except Exception:
            pass
    ln = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        (SLIDE_W - Inches(2.5)) / 2, Inches(3.3),
        Inches(2.5), Emu(20000)
    )
    ln.fill.solid()
    ln.fill.fore_color.rgb = COPPER
    ln.line.fill.background()
    add_text(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.8),
             spec.get("title", "Thank You"),
             size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.5),
             spec.get("subtitle", ""),
             size=18, color=COPPER, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.5),
             spec.get("brand", "Thorius Academy"),
             size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
             spec.get("contact", "www.thorius.com.tr"),
             size=14, color=GREY_LITE, align=PP_ALIGN.CENTER)


# ============================================================
# DISPATCH
# ============================================================
LAYOUT_FUNCS = {
    "cover": layout_cover,
    "section_divider": layout_section_divider,
    "content_typography": layout_content_typography,
    "three_columns": layout_three_columns,
    "timeline": layout_timeline,
    "bullets": layout_bullets,
    "closing": layout_closing,
}


# ============================================================
# MAIN ENTRY
# ============================================================
def build_presentation_impl(
    course_slug: str,
    section_id: str,
    slides: list,
    total_slides_in_course: int = 67,
    starting_page_no: int = 1,
) -> str:
    sec_dir = ensure_section(course_slug, section_id)
    out_path = sec_dir / "slides.pptx"

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    BLANK = prs.slide_layouts[6]

    page_no = starting_page_no
    for i, spec in enumerate(slides):
        layout = spec.get("layout", "bullets")
        s = prs.slides.add_slide(BLANK)

        if layout in ("cover", "section_divider", "closing"):
            # Full-bleed slides — no header/footer/bg pattern
            LAYOUT_FUNCS[layout](s, spec)
        else:
            # Content slide — header + body bg + footer
            add_body_bg(s)
            add_header(
                s,
                spec.get("title", ""),
                kicker=spec.get("kicker", None)
            )
            add_footer(s, page_no, total_slides_in_course)
            fn = LAYOUT_FUNCS.get(layout)
            if fn is None:
                add_text(s, Inches(1), Inches(3), Inches(11), Inches(1),
                         f"[Unknown layout: {layout}]",
                         size=24, color=COPPER, align=PP_ALIGN.CENTER)
            else:
                fn(s, spec)

        page_no += 1

    prs.save(str(out_path))
    mark_step(course_slug, section_id, "pptx", str(out_path))
    return (
        f"✓ Built {len(slides)} slide(s) → {out_path}\n"
        f"  Pages {starting_page_no}–{page_no - 1} of {total_slides_in_course}"
    )
